from __future__ import annotations

from copy import deepcopy
import io
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

TOKEN_PATTERN = re.compile(r"\{\{\s*([a-zA-Z0-9_\-]+)\s*\}\}")
IMAGE_TOKEN_PATTERN = re.compile(r"\{\{\s*image\s*:\s*([a-zA-Z0-9_\-]+)\s*\}\}")


@dataclass
class TemplateScanResult:
    text_tokens: List[str]
    image_tokens: List[str]
    slide_count: int


class PptTemplateEngine:
    """Token based PPT template engine for style-consistent generation.

    Text token format: {{field_name}}
    Image token format: {{image:field_name}}
    """

    def __init__(self, template_bytes: bytes) -> None:
        self._template_bytes = template_bytes

    def scan_tokens(self) -> TemplateScanResult:
        presentation = Presentation(io.BytesIO(self._template_bytes))
        text_tokens = set()
        image_tokens = set()

        for slide in presentation.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = shape.text or ""
                for match in TOKEN_PATTERN.findall(text):
                    text_tokens.add(match)
                for match in IMAGE_TOKEN_PATTERN.findall(text):
                    image_tokens.add(match)

        return TemplateScanResult(
            text_tokens=sorted(text_tokens),
            image_tokens=sorted(image_tokens),
            slide_count=len(presentation.slides),
        )

    def render(
        self,
        text_values: Dict[str, str],
        image_paths: Dict[str, Path],
    ) -> Tuple[bytes, List[str]]:
        """Render ppt from tokens and return (ppt_bytes, warnings)."""

        presentation = Presentation(io.BytesIO(self._template_bytes))
        warnings: List[str] = []

        for slide_idx, slide in enumerate(presentation.slides, start=1):
            to_replace_images = []
            missing_text_tokens = set()
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue

                original = shape.text or ""
                if not original:
                    continue

                image_match = IMAGE_TOKEN_PATTERN.fullmatch(original.strip())
                if image_match:
                    token = image_match.group(1)
                    image_path = image_paths.get(token)
                    if image_path and image_path.exists():
                        to_replace_images.append((shape, image_path, token))
                    else:
                        warnings.append(f"第 {slide_idx} 页缺少图片字段: {token}")
                    continue

                self._replace_text_tokens_in_place(
                    shape=shape,
                    text_values=text_values,
                    missing_text_tokens=missing_text_tokens,
                )

            for token in sorted(missing_text_tokens):
                warnings.append(f"第 {slide_idx} 页缺少文本字段: {token}")

            for shape, image_path, token in to_replace_images:
                self._replace_text_shape_with_image(slide, shape, image_path, token)

        output = io.BytesIO()
        presentation.save(output)
        return output.getvalue(), warnings

    @staticmethod
    def _replace_text_shape_with_image(slide, shape, image_path: Path, token: str) -> None:
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height

        element = shape._element
        element.getparent().remove(element)

        if "logo" in token.lower():
            picture_left, picture_top, picture_width, picture_height = PptTemplateEngine._calculate_contain_picture_box(
                image_path,
                left,
                top,
                width,
                height,
            )
            picture_left, picture_top, picture_width, picture_height = PptTemplateEngine._scale_picture_box_from_center(
                picture_left,
                picture_top,
                picture_width,
                picture_height,
                scale=2.0,
            )
            slide.shapes.add_picture(
                str(image_path),
                picture_left,
                picture_top,
                width=picture_width,
                height=picture_height,
            )
            return

        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

    @staticmethod
    def _calculate_contain_picture_box(image_path: Path, left: int, top: int, width: int, height: int) -> Tuple[int, int, int, int]:
        with Image.open(image_path) as image:
            image_width, image_height = image.size

        if image_width <= 0 or image_height <= 0:
            return left, top, width, height

        scale = min(width / image_width, height / image_height)
        picture_width = int(image_width * scale)
        picture_height = int(image_height * scale)
        picture_left = int(left + (width - picture_width) / 2)
        picture_top = int(top + (height - picture_height) / 2)
        return picture_left, picture_top, picture_width, picture_height

    @staticmethod
    def _scale_picture_box_from_center(left: int, top: int, width: int, height: int, scale: float) -> Tuple[int, int, int, int]:
        scaled_width = int(width * scale)
        scaled_height = int(height * scale)
        scaled_left = int(left - (scaled_width - width) / 2)
        scaled_top = int(top - (scaled_height - height) / 2)
        return scaled_left, scaled_top, scaled_width, scaled_height

    @staticmethod
    def _replace_text_tokens_in_place(shape, text_values: Dict[str, str], missing_text_tokens: set) -> None:
        """Replace tokens in each run to preserve original PPT formatting.

        This keeps bullet style, colors, spacing and paragraph-level settings from the template.
        """

        paragraphs = list(shape.text_frame.paragraphs)
        for paragraph in paragraphs:
            paragraph_text = paragraph.text or ""
            paragraph_token_matches = TOKEN_PATTERN.findall(paragraph_text)

            if len(paragraph_token_matches) == 1 and TOKEN_PATTERN.fullmatch(paragraph_text.strip()):
                token = paragraph_token_matches[0]
                replacement = text_values.get(token)
                if replacement is None:
                    missing_text_tokens.add(token)
                    replacement = ""
                normalized_replacement = PptTemplateEngine._normalize_replacement_text(replacement)
                replacement_lines = [line.strip() for line in normalized_replacement.split("\n") if line.strip()]
                if replacement_lines:
                    PptTemplateEngine._replace_paragraph_with_lines(paragraph, replacement_lines)
                else:
                    paragraph_style = PptTemplateEngine._capture_paragraph_font_style(paragraph)
                    paragraph.text = ""
                    PptTemplateEngine._apply_paragraph_font_style(paragraph, paragraph_style)
                continue

            if paragraph_token_matches:
                run_level_tokens = []
                multiline_replacement = False
                for run in paragraph.runs:
                    run_level_tokens.extend(TOKEN_PATTERN.findall(run.text or ""))

                for token in paragraph_token_matches:
                    replacement = text_values.get(token, "")
                    if "\n" in PptTemplateEngine._normalize_replacement_text(replacement):
                        multiline_replacement = True
                        break

                # Fallback to paragraph-level replacement when a token is split across runs
                # or replacement itself contains line breaks.
                if sorted(run_level_tokens) != sorted(paragraph_token_matches) or multiline_replacement:
                    replaced_text = paragraph_text
                    for token in paragraph_token_matches:
                        replacement = text_values.get(token)
                        if replacement is None:
                            missing_text_tokens.add(token)
                            replacement = ""
                        replacement = PptTemplateEngine._normalize_replacement_text(replacement)
                        replaced_text = re.sub(
                            rf"\{{\{{\s*{re.escape(token)}\s*\}}\}}",
                            replacement,
                            replaced_text,
                        )

                    if replaced_text != paragraph_text:
                        paragraph_style = PptTemplateEngine._capture_paragraph_font_style(paragraph)
                        paragraph.text = replaced_text
                        PptTemplateEngine._apply_paragraph_font_style(paragraph, paragraph_style)
                    continue

            if not paragraph.runs:
                replaced_text = paragraph_text
                for token in TOKEN_PATTERN.findall(paragraph_text):
                    replacement = text_values.get(token)
                    if replacement is None:
                        missing_text_tokens.add(token)
                        replacement = ""
                    replacement = PptTemplateEngine._normalize_replacement_text(replacement)
                    replaced_text = re.sub(
                        rf"\{{\{{\s*{re.escape(token)}\s*\}}\}}",
                        replacement,
                        replaced_text,
                    )
                if replaced_text != paragraph_text:
                    paragraph_style = PptTemplateEngine._capture_paragraph_font_style(paragraph)
                    paragraph.text = replaced_text
                    PptTemplateEngine._apply_paragraph_font_style(paragraph, paragraph_style)
                continue

            for run in paragraph.runs:
                run_text = run.text or ""
                if not run_text:
                    continue

                replaced_text = run_text
                for token in TOKEN_PATTERN.findall(run_text):
                    replacement = text_values.get(token)
                    if replacement is None:
                        missing_text_tokens.add(token)
                        replacement = ""
                    replacement = PptTemplateEngine._normalize_replacement_text(replacement)
                    replaced_text = re.sub(
                        rf"\{{\{{\s*{re.escape(token)}\s*\}}\}}",
                        replacement,
                        replaced_text,
                    )

                if replaced_text != run_text:
                    run.text = replaced_text

    @staticmethod
    def _replace_paragraph_with_lines(paragraph, lines: List[str]) -> None:
        """Replace one token paragraph with multiple paragraphs preserving template paragraph style."""

        paragraph_style = PptTemplateEngine._capture_paragraph_font_style(paragraph)
        source_xml = deepcopy(paragraph._p)
        paragraph.text = lines[0]
        PptTemplateEngine._apply_paragraph_font_style(paragraph, paragraph_style)
        source_p = paragraph._p

        for line in lines[1:]:
            text_frame = paragraph._parent
            new_paragraph = text_frame.add_paragraph()
            PptTemplateEngine._copy_paragraph_properties_xml(source_xml, new_paragraph._p)
            new_paragraph.text = line
            PptTemplateEngine._apply_paragraph_font_style(new_paragraph, paragraph_style)
            parent_element = new_paragraph._p.getparent()
            parent_element.remove(new_paragraph._p)
            source_p.addnext(new_paragraph._p)
            source_p = new_paragraph._p

    @staticmethod
    def _copy_paragraph_properties_xml(source_p, target_p) -> None:
        source_ppr = None
        for child in list(source_p):
            if child.tag.endswith("}pPr"):
                source_ppr = deepcopy(child)
                break

        for child in list(target_p):
            if child.tag.endswith("}pPr"):
                target_p.remove(child)

        if source_ppr is not None:
            target_p.insert(0, source_ppr)

    @staticmethod
    def _capture_paragraph_font_style(paragraph) -> Dict[str, object]:
        run = paragraph.runs[0] if paragraph.runs else None
        font = run.font if run else paragraph.font
        return PptTemplateEngine._capture_font_style(font)

    @staticmethod
    def _apply_paragraph_font_style(paragraph, style: Dict[str, object]) -> None:
        for run in paragraph.runs:
            PptTemplateEngine._apply_font_style(run.font, style)

    @staticmethod
    def _normalize_replacement_text(value: str) -> str:
        """Normalize replacement text and keep line breaks visible in one paragraph.

        For run-level replacement, use `\n` to avoid exposing Office escape markers
        like `_x000B_` as plain text.
        """

        normalized = (
            value.replace("_x000D_", "\n")
            .replace("x000D", "\n")
            .replace("_x000B_", "\n")
            .replace("x000B", "\n")
            .replace("\r\n", "\n")
            .replace("\r", "\n")
        )
        return normalized

    @staticmethod
    def _capture_shape_text_style(
        shape,
        prefer_token_style: bool = False,
    ) -> Tuple[Optional[object], Optional[object], Dict[str, object]]:
        text_frame = shape.text_frame
        if not text_frame.paragraphs:
            return None, None, {}

        paragraph = text_frame.paragraphs[0]
        run = paragraph.runs[0] if paragraph.runs else None

        if prefer_token_style:
            for p in text_frame.paragraphs:
                for r in p.runs:
                    if TOKEN_PATTERN.search(r.text or ""):
                        paragraph = p
                        run = r
                        break
                else:
                    continue
                break

        font = run.font if run else paragraph.font
        style = PptTemplateEngine._capture_font_style(font)
        return paragraph.alignment, paragraph.level, style

    @staticmethod
    def _set_shape_text_preserve_style(
        shape,
        text: str,
        alignment: Optional[object],
        level: Optional[object],
        style: Dict[str, object],
    ) -> None:
        text_frame = shape.text_frame
        text_frame.clear()

        lines = text.splitlines() or [""]
        paragraph = text_frame.paragraphs[0]
        paragraph.text = lines[0]

        for extra_line in lines[1:]:
            p = text_frame.add_paragraph()
            p.text = extra_line

        for p in text_frame.paragraphs:
            if alignment is not None:
                p.alignment = alignment
            if level is not None:
                p.level = level

            for r in p.runs:
                PptTemplateEngine._apply_font_style(r.font, style)

    @staticmethod
    def _capture_font_style(font) -> Dict[str, object]:
        color = font.color
        style: Dict[str, object] = {
            "size": font.size,
            "name": font.name,
            "bold": font.bold,
            "italic": font.italic,
            "underline": font.underline,
            "color_type": getattr(color, "type", None),
            "color_rgb": None,
            "color_theme": None,
            "color_brightness": None,
        }

        try:
            style["color_rgb"] = color.rgb
        except Exception:
            style["color_rgb"] = None

        try:
            style["color_theme"] = color.theme_color
        except Exception:
            style["color_theme"] = None

        try:
            style["color_brightness"] = color.brightness
        except Exception:
            style["color_brightness"] = None

        return style

    @staticmethod
    def _apply_font_style(font, style: Dict[str, object]) -> None:
        font.size = style.get("size")
        font.name = style.get("name")
        font.bold = style.get("bold")
        font.italic = style.get("italic")
        font.underline = style.get("underline")

        color_rgb = style.get("color_rgb")
        color_theme = style.get("color_theme")
        color_brightness = style.get("color_brightness")

        if color_rgb is not None:
            try:
                font.color.rgb = color_rgb
            except Exception:
                pass
        elif color_theme is not None:
            try:
                font.color.theme_color = color_theme
                if color_brightness is not None:
                    font.color.brightness = color_brightness
            except Exception:
                pass
